# server/ppt_analyzer.py
# This script performs the core AI logic for matching test cases to PowerPoints.
# It's called by the Node.js server and communicates over stdin/stdout.
# This version includes confidence levels, duplicate detection, robust validation, and performance optimizations.

import sys
import os
import pandas as pd
from pptx import Presentation
from sentence_transformers import SentenceTransformer, util
import torch
import json
import io

# --- Centralized configuration for thresholds ---
# This makes it easier to adjust the matching sensitivity in one place.
SIMILARITY_CONFIDENCE_THRESHOLDS = {
    "high": 0.7,
    "medium": 0.5,
    "low": 0.3
}
DUPLICATE_DESCRIPTION_THRESHOLD = 0.85

def read_ppt_content(ppt_path):
    """
    Extracts all text from a PowerPoint presentation, handling potential errors gracefully.
    """
    try:
        prs = Presentation(ppt_path)
        text_runs = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        text_runs.append(run.text)
        return " ".join(text_runs)
    except Exception as e:
        # Log a warning to stderr so Node.js can see it, but don't crash the script.
        print(f"Warning: Could not read {os.path.basename(ppt_path)}. Error: {e}", file=sys.stderr)
        return ""

def get_confidence_level(score):
    """
    Categorizes a similarity score into a human-readable confidence level.
    """
    if score >= SIMILARITY_CONFIDENCE_THRESHOLDS["high"]:
        return "High"
    elif score >= SIMILARITY_CONFIDENCE_THRESHOLDS["medium"]:
        return "Medium"
    elif score >= SIMILARITY_CONFIDENCE_THRESHOLDS["low"]:
        return "Low"
    else:
        return "No confident match"

def main(knowledge_repo_path, excel_buffer):
    """
    Main function to perform the analysis.
    - Loads a pre-trained sentence transformer model.
    - Reads all PPTs from the knowledge repository.
    - Validates and reads the user-uploaded Excel file from a buffer.
    - Encodes text from both sources into vector embeddings.
    - Computes similarity and finds the best match for each test case.
    - Handles duplicates, adds confidence levels, and formats the output as JSON.
    """
    # Load the sentence transformer model. This might download the model on the first run.
    model = SentenceTransformer('all-MiniLM-L6-v2')

    # 1. Process all PowerPoints from the knowledge repository folder.
    ppt_contents = {}
    for filename in os.listdir(knowledge_repo_path):
        if filename.endswith(".pptx"):
            ppt_path = os.path.join(knowledge_repo_path, filename)
            content = read_ppt_content(ppt_path)
            if content:
                ppt_contents[filename] = content
    
    if not ppt_contents:
        print(json.dumps({"error": "No readable PowerPoint files found in the repository."}), file=sys.stdout)
        sys.exit(1)

    # 2. Read and robustly validate the uploaded Excel file.
    try:
        df = pd.read_excel(excel_buffer, header=0)
        # Clean column names to be case-insensitive and ignore whitespace.
        df.columns = [col.strip() for col in df.columns]
        required_cols = ['Test Case ID', 'Test Case Description']
        
        if not all(col in df.columns for col in required_cols):
            missing = ", ".join([col for col in required_cols if col not in df.columns])
            print(json.dumps({"error": f"Excel file is missing required columns: {missing}"}), file=sys.stdout)
            sys.exit(1)

        # Ensure we only work with the columns we need.
        df = df[required_cols]
        df.dropna(subset=['Test Case Description'], inplace=True)
        if df.empty:
            print(json.dumps({"error": "No valid test case descriptions found in the Excel file."}), file=sys.stdout)
            sys.exit(1)
        test_cases = df['Test Case Description'].tolist()
    except Exception as e:
        print(json.dumps({"error": f"Failed to read or validate the Excel file. Details: {e}"}), file=sys.stdout)
        sys.exit(1)

    # 3. Create embeddings using torch.no_grad() for a performance boost.
    # This tells PyTorch not to calculate gradients, which is unnecessary for inference.
    with torch.no_grad():
        ppt_filenames = list(ppt_contents.keys())
        ppt_texts = list(ppt_contents.values())
        ppt_embeddings = model.encode(ppt_texts, convert_to_tensor=True)
        test_case_embeddings = model.encode(test_cases, convert_to_tensor=True)

        # 4. Compute cosine similarity between all test cases and all PowerPoints.
        cosine_scores = util.cos_sim(test_case_embeddings, ppt_embeddings)

    # 5. Find the best match for each test case and apply enhanced logic.
    results = []
    first_match_for_ppt = {}

    for i, row in df.iterrows():
        test_case_id = row['Test Case ID']
        test_case_desc = row['Test Case Description']
        
        scores_for_case = cosine_scores[i]
        top_result_index = torch.argmax(scores_for_case).item()
        
        best_ppt_filename = ppt_filenames[top_result_index]
        similarity_score = scores_for_case[top_result_index].item()
        confidence = get_confidence_level(similarity_score)
        
        note = ""
        final_ppt_match = best_ppt_filename

        # Handle low confidence matches by providing a clear message.
        if confidence == "No confident match":
            final_ppt_match = "No relevant deck found"
            note = "The top match score was too low for a confident recommendation."
        # Handle duplicate detection for confident matches.
        elif best_ppt_filename in first_match_for_ppt:
            original_match_info = first_match_for_ppt[best_ppt_filename]
            original_match_desc = original_match_info['description']
            
            # To check for duplication, compare the current test case description
            # with the description of the *first* test case that matched this PPT.
            with torch.no_grad():
                desc_embeddings = model.encode([test_case_desc, original_match_desc], convert_to_tensor=True)
                desc_similarity = util.cos_sim(desc_embeddings[0], desc_embeddings[1]).item()
            
            # If descriptions are very similar, flag as a potential duplicate.
            if desc_similarity > DUPLICATE_DESCRIPTION_THRESHOLD:
                 note = f"Semantically similar to Test Case ID '{original_match_info['id']}'. Same PPT recommended."
        else:
            # This is the first time we've seen a confident match for this PPT. Record it.
            first_match_for_ppt[best_ppt_filename] = {'id': test_case_id, 'description': test_case_desc}

        results.append({
            "Test Case ID": test_case_id,
            "Test Case Description": test_case_desc,
            "Matched PPT": final_ppt_match,
            "Similarity Score": f"{similarity_score:.2f}",
            "Confidence Level": confidence,
            "Note": note
        })

    # 6. Output results as a JSON string to stdout, which Node.js will capture.
    print(json.dumps(results, indent=4))

if __name__ == "__main__":
    if len(sys.argv) != 2:
        print("Usage: python ppt_analyzer.py <knowledge_repo_path>", file=sys.stderr)
        sys.exit(1)
        
    knowledge_repo_path = sys.argv[1]
    # Read the Excel file content from standard input as a binary buffer.
    excel_buffer = io.BytesIO(sys.stdin.buffer.read())
    
    main(knowledge_repo_path, excel_buffer)
